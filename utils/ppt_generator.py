from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from typing import List, Dict
import os


def _add_textbox(slide, left, top, width, height, text, font_size=12, bold=False, color=(0, 0, 0)):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    run = paragraph.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    return textbox


def _add_bullet_list(slide, left, top, width, height, title, items, title_size=12, item_size=11):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.clear()

    title_paragraph = frame.paragraphs[0]
    title_paragraph.text = title
    title_paragraph.font.size = Pt(title_size)
    title_paragraph.font.bold = True

    for item in items:
        p = frame.add_paragraph()
        p.text = f"• {item}"
        p.level = 1
        p.font.size = Pt(item_size)
    return textbox


def _add_picture(slide, image_path: str, left, top, width=None, height=None):
    if not os.path.exists(image_path):
        return None
    if width and height:
        return slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    return slide.shapes.add_picture(image_path, left, top)


def generate_region_dashboard_ppt(region_cards: List[Dict], output_path: str) -> str:
    """Generate a PowerPoint dashboard summary for one or more regions."""
    prs = Presentation()

    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]

    for card in region_cards:
        region = card.get("region", "Region")
        date_range = card.get("date_range", "")
        metrics = card.get("metrics", {})
        top_stations = card.get("top_stations", [])
        top_feeders = card.get("top_feeders", [])
        chart_image = card.get("chart_image")

        slide = prs.slides.add_slide(blank_layout)
        _add_textbox(slide, Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.6), f"{region} Dashboard", font_size=30, bold=True)
        _add_textbox(slide, Inches(0.4), Inches(0.9), Inches(9.2), Inches(0.4), date_range, font_size=14)

        metrics_text = "\n".join([f"{label}: {value}" for label, value in metrics.items()])
        _add_textbox(slide, Inches(0.4), Inches(1.5), Inches(4.7), Inches(3.0), metrics_text, font_size=12)

        if chart_image:
            _add_picture(slide, chart_image, Inches(5.2), Inches(1.5), width=Inches(4.0), height=Inches(3.0))

        if top_stations:
            station_items = [f"{row[0]}: {row[1]} hrs, {row[2]} MWh" for row in top_stations]
            _add_bullet_list(slide, Inches(0.4), Inches(4.6), Inches(4.7), Inches(2.4), "Top Stations", station_items[:6])

        if top_feeders:
            feeder_items = [f"{row[0]}: {row[1]} hrs, {row[2]} MWh" for row in top_feeders]
            _add_bullet_list(slide, Inches(5.2), Inches(4.6), Inches(4.7), Inches(2.4), "Top Feeders", feeder_items[:6])

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path
